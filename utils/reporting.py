# utils/reporting.py
import pandas as pd
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches

def save_metrics_rows(rows, out_path="outputs/reports.xlsx"):
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    df = pd.DataFrame(rows)
    # For counts dict, convert to string for Excel readability
    if 'counts' in df.columns:
        df['counts'] = df['counts'].apply(lambda x: str(x))
    df.to_excel(out_path, index=False)
    return out_path

def create_brief_ppt(samples, out_ppt="outputs/demo_slides.pptx"):
    prs = Presentation()
    for s in samples:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = s.get('title', '')
        left = Inches(0.5); top = Inches(1.2)
        if s.get('before_path') and s.get('after_path'):
            slide.shapes.add_picture(s['before_path'], left, top, width=Inches(4.0))
            slide.shapes.add_picture(s['after_path'], left+Inches(4.2), top, width=Inches(4.0))
    os.makedirs(os.path.dirname(out_ppt), exist_ok=True)
    prs.save(out_ppt)
    return out_ppt